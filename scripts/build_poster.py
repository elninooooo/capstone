"""Populate the capstone poster PPTX with real content, figures, and tables.

Target: 1234803-W02-A Traffic-Optimized Data Sharing Framework for IoT
Devices using Smart Contracts and Distributed Storage.pptx (36" x 24").
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy

ROOT = Path(__file__).resolve().parent.parent
PPTX_NAME = (
    "1234803-W02-A Traffic-Optimized Data Sharing Framework for IoT "
    "Devices using Smart Contracts and Distributed Storage.pptx"
)
PPTX_PATH = ROOT / PPTX_NAME
FIG_DIR = ROOT / "figures_testing"

NAVY = RGBColor(0x1F, 0x3A, 0x68)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xC9, 0x53, 0x0F)

FULL_TITLE = (
    "A Traffic-Optimized Data Sharing Framework for IoT Devices "
    "using Smart Contracts and Distributed Storage"
)


def _clear_text_frame(tf):
    """Remove every paragraph in a text frame, leaving a single empty one."""
    txBody = tf._txBody
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)
    # python-pptx expects at least one paragraph; add a fresh empty one.
    from pptx.oxml.ns import nsmap  # noqa: F401

    # Use the tf.add_paragraph API by re-accessing via tf.paragraphs trick:
    from lxml import etree

    p = etree.SubElement(txBody, qn("a:p"))
    return p


def _add_paragraph(tf, first=False):
    """Return a blank new paragraph on the text frame (keeping old ones)."""
    if first and len(tf.paragraphs) == 1 and not tf.paragraphs[0].text:
        return tf.paragraphs[0]
    return tf.add_paragraph()


def _write_runs(para, runs, align=PP_ALIGN.LEFT, space_after=4, line_spacing=1.05):
    """`runs` is a list of (text, size_pt, bold, color, italic?) tuples."""
    para.alignment = align
    para.space_after = Pt(space_after)
    para.line_spacing = line_spacing
    for item in runs:
        text = item[0]
        size = item[1] if len(item) > 1 else 24
        bold = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else DARK
        italic = item[4] if len(item) > 4 else False
        r = para.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Calibri"


def _write_lines(tf, lines):
    """`lines` is list of dicts: {runs, align, space_after, bullet}."""
    # Clear all existing paragraphs
    from lxml import etree

    txBody = tf._txBody
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)

    for i, line in enumerate(lines):
        p_el = etree.SubElement(txBody, qn("a:p"))
        # pPr for alignment & spacing
        pPr = etree.SubElement(p_el, qn("a:pPr"))
        align = line.get("align", "l")
        pPr.set("algn", {"l": "l", "c": "ctr", "r": "r"}[align])
        if "space_before" in line:
            pPr.set("spcBef", "")
        for r in line.get("runs", []):
            run_el = etree.SubElement(p_el, qn("a:r"))
            rPr = etree.SubElement(run_el, qn("a:rPr"))
            rPr.set("lang", "en-US")
            rPr.set("sz", str(int(r.get("size", 24) * 100)))
            if r.get("bold"):
                rPr.set("b", "1")
            if r.get("italic"):
                rPr.set("i", "1")
            color = r.get("color", DARK)
            solidFill = etree.SubElement(rPr, qn("a:solidFill"))
            srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgb.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))
            latin = etree.SubElement(rPr, qn("a:latin"))
            latin.set("typeface", "Calibri")
            t = etree.SubElement(run_el, qn("a:t"))
            t.text = r["text"]
        if not line.get("runs"):
            # empty paragraph: still needs endParaRPr
            endPr = etree.SubElement(p_el, qn("a:endParaRPr"))
            endPr.set("lang", "en-US")


# ---------------------------------------------------------------------------


def build_intro(shape):
    """Left upper column — Motivation, Gap, Objectives."""
    tf = shape.text_frame
    tf.word_wrap = True
    lines = []

    lines.append({"runs": [{"text": "Motivation", "size": 28, "bold": True, "color": NAVY}]})
    lines.append(
        {
            "runs": [
                {
                    "text": (
                        "Explosive IoT growth produces massive sensor data. "
                        "Centralised clouds suffer from single-point failure, "
                        "privacy leakage, tampering risk and high cost."
                    ),
                    "size": 24,
                }
            ]
        }
    )
    lines.append({"runs": []})

    lines.append({"runs": [{"text": "Research Gap", "size": 28, "bold": True, "color": NAVY}]})
    lines.append(
        {
            "runs": [
                {
                    "text": (
                        "Prior blockchain-IoT designs either store every "
                        "payload on-chain (expensive, congesting the network) "
                        "or only log a hash without role-based access control "
                        "and end-to-end traffic optimisation."
                    ),
                    "size": 24,
                }
            ]
        }
    )
    lines.append({"runs": []})

    lines.append({"runs": [{"text": "Objectives", "size": 28, "bold": True, "color": NAVY}]})
    for obj in [
        "Hybrid on-chain + off-chain storage to minimise on-chain load.",
        "Smart contract governed, decentralised role-based access control.",
        "Batch uploading to cut transaction count and gas cost.",
        "Hash-based integrity verification with a full web dashboard.",
    ]:
        lines.append(
            {
                "runs": [
                    {"text": "•  ", "size": 24, "bold": True, "color": NAVY},
                    {"text": obj, "size": 24},
                ]
            }
        )
    _write_lines(tf, lines)


def build_methods(shape):
    """Left lower column — stack + key components (architecture figure is a separate picture)."""
    tf = shape.text_frame
    tf.word_wrap = True
    lines = []

    lines.append(
        {"runs": [{"text": "Technology Stack", "size": 26, "bold": True, "color": NAVY}]}
    )
    lines.append(
        {
            "runs": [
                {
                    "text": (
                        "Solidity 0.8.20 · Hardhat 2.28 · Flask 3 · "
                        "Web3.py 7 · React 19 + Ant Design 6 · "
                        "Python IoT Simulator · Docker"
                    ),
                    "size": 22,
                }
            ]
        }
    )
    lines.append({"runs": []})

    lines.append(
        {"runs": [{"text": "Key Components", "size": 26, "bold": True, "color": NAVY}]}
    )
    for label, desc in [
        ("Smart Contract", "device registry, ACL and integrity proof."),
        ("Off-chain Store", "content-addressable (IPFS-style CID)."),
        ("Flask Backend", "16 REST endpoints + Web3.py bridge."),
        ("React Dashboard", "5 routes + Python batch IoT simulator."),
    ]:
        lines.append(
            {
                "runs": [
                    {"text": "•  ", "size": 22, "bold": True, "color": NAVY},
                    {"text": label + ": ", "size": 22, "bold": True},
                    {"text": desc, "size": 22},
                ]
            }
        )
    _write_lines(tf, lines)


def build_results(shape):
    """Middle column headline outcomes (figures & tables added separately)."""
    tf = shape.text_frame
    tf.word_wrap = True
    lines = []

    lines.append(
        {
            "align": "c",
            "runs": [
                {"text": "Headline Outcomes", "size": 28, "bold": True, "color": NAVY},
            ],
        }
    )
    lines.append(
        {
            "align": "c",
            "runs": [
                {"text": "75-85%", "size": 56, "bold": True, "color": ACCENT},
                {"text": "   on-chain payload reduction", "size": 20, "bold": True},
            ],
        }
    )
    lines.append(
        {
            "align": "c",
            "runs": [
                {"text": "87.5%", "size": 56, "bold": True, "color": ACCENT},
                {"text": "   fewer transactions at batch size 8", "size": 20, "bold": True},
            ],
        }
    )
    lines.append(
        {
            "align": "c",
            "runs": [
                {"text": "27 / 27", "size": 56, "bold": True, "color": ACCENT},
                {"text": "   smart-contract unit tests passed", "size": 20, "bold": True},
            ],
        }
    )

    _write_lines(tf, lines)


def build_conclusions(shape):
    """Right upper column — Achievements + Future Work."""
    tf = shape.text_frame
    tf.word_wrap = True
    lines = []

    lines.append(
        {"runs": [{"text": "Achievements", "size": 28, "bold": True, "color": NAVY}]}
    )
    for a in [
        "Full-stack prototype (contract + backend + frontend + simulator) built and tested.",
        "75-85% on-chain payload compressed; 87.5% fewer transactions at batch size 8.",
        "RBAC cross-device isolation: all 9 access scenarios pass.",
        "Tamper detection fires on any payload hash mismatch.",
        "16 REST endpoints and 5 frontend routes fully integrated.",
    ]:
        lines.append(
            {
                "runs": [
                    {"text": "•  ", "size": 22, "bold": True, "color": NAVY},
                    {"text": a, "size": 22},
                ]
            }
        )

    lines.append({"runs": []})
    lines.append(
        {"runs": [{"text": "Future Work", "size": 28, "bold": True, "color": NAVY}]}
    )
    for f in [
        "Migrate off-chain layer to a production IPFS cluster.",
        "Deploy on Layer-2 (Polygon / Arbitrum) to benchmark real gas cost.",
        "Extend RBAC to attribute-based access with an encryption layer.",
    ]:
        lines.append(
            {
                "runs": [
                    {"text": "•  ", "size": 22, "bold": True, "color": NAVY},
                    {"text": f, "size": 22},
                ]
            }
        )
    _write_lines(tf, lines)


def build_references(shape):
    """Right lower column — 7 concise references."""
    tf = shape.text_frame
    tf.word_wrap = True
    refs = [
        "Nakamoto S. (2008). Bitcoin: A Peer-to-Peer Electronic Cash System.",
        "Wood G. (2014). Ethereum: A Secure Decentralised Generalised Transaction Ledger.",
        "Benet J. (2014). IPFS - Content-Addressed, Versioned, P2P File System.",
        "Christidis K. & Devetsikiotis M. (2016). Blockchains and Smart Contracts for IoT. IEEE Access 4.",
        "Zhang Y. et al. (2019). Smart Contract-Based Access Control for the IoT. IEEE IoT-J 6(2).",
        "Zheng Q. et al. (2018). An IPFS-Blockchain Cooperative Architecture. IEEE BigData Congress.",
        "Buterin V. (2014). Ethereum White Paper.",
    ]
    lines = []
    for i, r in enumerate(refs, start=1):
        lines.append(
            {
                "runs": [
                    {"text": f"[{i}] ", "size": 18, "bold": True, "color": NAVY},
                    {"text": r, "size": 18, "color": GREY},
                ]
            }
        )
    _write_lines(tf, lines)


# ---------------------------------------------------------------------------


def add_table(slide, left, top, width, height, header, rows, header_fill=NAVY,
              header_text_size=18, body_text_size=18):
    shape = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(header),
        left=left,
        top=top,
        width=width,
        height=height,
    )
    tbl = shape.table
    # Header row
    for ci, h in enumerate(header):
        cell = tbl.cell(0, ci)
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        _write_lines(
            tf,
            [
                {
                    "align": "c",
                    "runs": [
                        {
                            "text": h,
                            "size": header_text_size,
                            "bold": True,
                            "color": RGBColor(0xFF, 0xFF, 0xFF),
                        }
                    ],
                }
            ],
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill

    # Body rows
    for ri, row in enumerate(rows, start=1):
        for ci, v in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            _write_lines(
                cell.text_frame,
                [
                    {
                        "align": "c",
                        "runs": [
                            {
                                "text": str(v),
                                "size": body_text_size,
                                "bold": ci == 0,
                                "color": DARK,
                            }
                        ],
                    }
                ],
            )
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF8)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shape


# ---------------------------------------------------------------------------


def main():
    prs = Presentation(str(PPTX_PATH))
    slide = prs.slides[0]

    # Map shapes by name (more robust than index)
    by_name = {sh.name: sh for sh in slide.shapes}
    shapes = list(slide.shapes)

    # --- Title: extend to full title ---
    title_shape = shapes[0]
    _write_lines(
        title_shape.text_frame,
        [
            {
                "align": "l",
                "runs": [{"text": FULL_TITLE, "size": 64, "bold": True, "color": NAVY}],
            }
        ],
    )

    # --- Column bodies ---
    intro_shape = shapes[3]       # TextBox 12
    results_shape = shapes[4]     # TextBox 15
    conclusions_shape = shapes[5] # TextBox 16
    methods_shape = shapes[6]     # TextBox 17
    references_shape = shapes[7]  # TextBox 18

    build_intro(intro_shape)
    build_methods(methods_shape)
    build_results(results_shape)
    build_conclusions(conclusions_shape)
    build_references(references_shape)

    # Re-position the Methods body text box so it sits BELOW the architecture
    # figure (the figure lives at T=14.15 with height 4.9" so text starts at 19.15").
    methods_shape.top = Inches(19.15)
    methods_shape.height = Inches(4.2)

    # The Results body text box goes BETWEEN the flow figure (ends T=10.15) and
    # Table A (starts T=14.9). Tighten to that slot.
    results_shape.top = Inches(10.3)
    results_shape.height = Inches(4.5)

    # --- Figures ---
    # 1) Architecture diagram in Methods column (top of the methods body).
    arch = FIG_DIR / "system_architecture.png"
    if arch.exists():
        slide.shapes.add_picture(
            str(arch),
            left=Inches(0.95),
            top=Inches(14.15),
            width=Inches(11.0),
            height=Inches(4.9),
        )

    # 2) Data-flow diagram at top of Results column.
    flow = FIG_DIR / "system_flowchart.png"
    if flow.exists():
        slide.shapes.add_picture(
            str(flow),
            left=Inches(12.7),
            top=Inches(5.05),
            width=Inches(11.25),
            height=Inches(5.1),
        )

    # 3) Test pyramid in lower Results column.
    pyramid = FIG_DIR / "T1_test_pyramid.png"
    if pyramid.exists():
        slide.shapes.add_picture(
            str(pyramid),
            left=Inches(12.7),
            top=Inches(18.4),
            width=Inches(5.3),
            height=Inches(4.6),
        )

    # --- Two compact tables in the Results column ---
    add_table(
        slide,
        left=Inches(12.8),
        top=Inches(14.9),
        width=Inches(11.1),
        height=Inches(1.6),
        header=["Scenario", "On-chain (bytes)", "Off-chain (bytes)"],
        rows=[
            ["Single record", "~120", "~450"],
            ["Batch of 8", "~150", "~3 200"],
            ["Reduction", "75-85 %", "-"],
        ],
        header_text_size=20,
        body_text_size=20,
    )

    add_table(
        slide,
        left=Inches(18.2),
        top=Inches(18.4),
        width=Inches(5.7),
        height=Inches(4.5),
        header=["Batch", "Tx / 8 recs", "Reduction"],
        rows=[
            ["1", "8", "baseline"],
            ["2", "4", "50 %"],
            ["4", "2", "75 %"],
            ["8", "1", "87.5 %"],
        ],
        header_text_size=20,
        body_text_size=20,
    )

    # --- Save ---
    prs.save(str(PPTX_PATH))
    print(f"Saved: {PPTX_PATH}")


if __name__ == "__main__":
    main()
